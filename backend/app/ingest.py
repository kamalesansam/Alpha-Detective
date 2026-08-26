"""Document lifecycle: validate -> store raw -> parse -> chunk -> (embed) -> commit.

Per-file failures are reported per file (HTTP stays 200); request-level
violations are the API layer's job. Heavy work runs in `asyncio.to_thread`
so /api/health and /api/documents stay responsive mid-index. Ingest and
delete are serialized behind a single asyncio lock — correctness beats
parallel ingest.

v1.2 adds five formats (.xlsx .pptx .html .htm .json), a per-doc `tables`
count, a chunk inventory read path, and startup auto-seeding. Every new
extraction path is bounded by a named cap (§2) and every cap failure is a
clean per-file `failed` entry — never a 500, never a partial commit.
"""

from __future__ import annotations

import asyncio
import csv
import hashlib
import io
import json
import logging
import re
import shutil
import uuid
import zipfile
from dataclasses import dataclass, field
from datetime import datetime, timezone
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, NamedTuple, Optional, Union

from . import config, providers, stores

logger = logging.getLogger("alpha.ingest")

ALLOWED_EXTS = (
    ".pdf", ".docx", ".txt", ".md", ".csv",
    ".xlsx", ".pptx", ".html", ".htm", ".json",
)
MAX_FILE_BYTES = 25 * 1024 * 1024
MAX_FILES_PER_REQUEST = 20
# Request-level ceiling for POST /api/documents: 20 files x 25 MB + multipart
# framing slack. Checked against Content-Length BEFORE any body buffering.
MAX_REQUEST_BYTES = MAX_FILES_PER_REQUEST * MAX_FILE_BYTES + 8 * 1024 * 1024
CSV_WINDOW_ROWS = 40

# --- extraction caps (§2) -----------------------------------------------------
# Every one of these is a *bound on work*, not a quality knob: an untrusted
# upload must never be able to make the parser allocate or spin without limit.
OOXML_MAX_ENTRIES = 5000
OOXML_MAX_UNCOMPRESSED_BYTES = 100 * 1024 * 1024
OOXML_MAX_COMPRESSION_RATIO = 200
XLSX_MAX_SHEETS = 50
XLSX_MAX_CELLS = 200_000
XLSX_WINDOW_ROWS = 40
PPTX_MAX_SLIDES = 500
PPTX_MAX_TABLE_CELLS = 20_000
# Format-agnostic ceiling on ACCUMULATED extracted text, enforced INCREMENTALLY
# during parsing for all ten extensions (it subsumes the old HTML-only cap at
# the same value). It must abort mid-parse: a length check after the fact is
# the very allocation it exists to prevent. ~5M chars is ~2,800 chunks — an
# embedding bill that would exhaust the free tier on a single upload.
MAX_EXTRACTED_TEXT_CHARS = 5_000_000
JSON_MAX_DEPTH = 20
JSON_MAX_NODES = 200_000
JSON_WINDOW_LINES = 40

# OOXML containers get the zip-bomb guard before any library opens them.
OOXML_EXTS = (".docx", ".xlsx", ".pptx")

UUID4_RE = re.compile(
    r"^[0-9a-f]{8}-[0-9a-f]{4}-4[0-9a-f]{3}-[89ab][0-9a-f]{3}-[0-9a-f]{12}\Z",
    re.IGNORECASE,  # \Z, never $: `$` also matches before a trailing newline,
)                   # and this regex is documented as the path-traversal boundary

CHUNK_PREVIEW_MAX_CHARS = 200  # §1.8 preview cap

_ingest_lock = asyncio.Lock()


class ExtractionCapExceeded(Exception):
    """A named extraction cap was hit. `message` is the verbatim §1.3 string.

    Always caught by `_ingest_one` and turned into a `failed` entry — it must
    never reach main.py's catch-all (that would be a 500 for what is a clean,
    expected, per-file refusal).
    """

    def __init__(self, message: str):
        super().__init__(message)
        self.message = message


class _TextBudget:
    """Accumulated extracted-text guard, checked as text is produced.

    Every parser adds each fragment BEFORE retaining it, so an untrusted
    document trips the cap while the parser is still streaming rather than
    after it has already materialised a gigabyte. Caps compose: whichever
    trips first wins.
    """

    __slots__ = ("used",)

    def __init__(self) -> None:
        self.used = 0

    def add(self, text) -> None:
        self.used += len(text) if not isinstance(text, int) else text
        if self.used > MAX_EXTRACTED_TEXT_CHARS:
            raise ExtractionCapExceeded(
                f"extracted text too large (cap: {MAX_EXTRACTED_TEXT_CHARS} characters)"
            )

    def keep(self, text: str) -> str:
        """add() + return the text, for use inline in an expression."""
        self.add(text)
        return text


# --- filename & content validation -------------------------------------------

def sanitize_filename(name: str) -> str:
    """Basename only; strip separators/NULs/control chars; cap 120; never empty."""
    base = (name or "").replace("\\", "/").rsplit("/", 1)[-1]
    base = "".join(ch for ch in base if ch not in ("\x00",) and (ord(ch) >= 32) and ch != "\x7f")
    base = base.strip().strip(".")
    if len(base) > 120:
        stem, dot, suffix = base.rpartition(".")
        if dot and 0 < len(suffix) <= 10:
            base = stem[: 120 - len(suffix) - 1].rstrip(".") + "." + suffix
        else:
            base = base[:120]
    return base or "file"


def sniff_ok(ext: str, head: bytes) -> bool:
    """Magic-byte / content sniff per CONTRACTS.md §1.3."""
    if ext == ".pdf":
        return head.startswith(b"%PDF")
    if ext in OOXML_EXTS:
        return head.startswith(b"PK\x03\x04")
    if ext in (".txt", ".md", ".csv", ".html", ".htm", ".json"):
        if b"\x00" in head:
            return False
        try:
            head.decode("utf-8")
        except UnicodeDecodeError:
            try:
                head.decode("latin-1")
            except UnicodeDecodeError:
                return False
        return True
    return False


def _decode_text(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1")


# --- OOXML container guard ----------------------------------------------------

_ZIP_BOMB_MESSAGE = "archive expands too much (possible zip bomb)"


def _guard_zip_infolist(infos: list) -> None:
    """The three OOXML caps, computed from the central directory only.

    No entry is ever extracted here: `ZipFile.infolist()` reads sizes that the
    archive declares, which is exactly what a bomb lies about in a way we can
    detect cheaply (ratio) and refuse before decompressing a byte.
    """
    if len(infos) > OOXML_MAX_ENTRIES:
        raise ExtractionCapExceeded(_ZIP_BOMB_MESSAGE)
    total_uncompressed = 0
    total_compressed = 0
    for info in infos:
        size = int(getattr(info, "file_size", 0) or 0)
        packed = int(getattr(info, "compress_size", 0) or 0)
        total_uncompressed += size
        total_compressed += packed
        if packed > 0 and size / packed > OOXML_MAX_COMPRESSION_RATIO:
            raise ExtractionCapExceeded(_ZIP_BOMB_MESSAGE)
    if total_uncompressed > OOXML_MAX_UNCOMPRESSED_BYTES:
        raise ExtractionCapExceeded(_ZIP_BOMB_MESSAGE)
    if total_compressed > 0 and total_uncompressed / total_compressed > OOXML_MAX_COMPRESSION_RATIO:
        raise ExtractionCapExceeded(_ZIP_BOMB_MESSAGE)


def ooxml_guard(path: Path) -> None:
    """Raise ExtractionCapExceeded when an OOXML file is a zip bomb.

    Runs for .docx/.xlsx/.pptx BEFORE docx/openpyxl/python-pptx touch the file.
    """
    with zipfile.ZipFile(str(path)) as zf:
        _guard_zip_infolist(zf.infolist())


def ooxml_guard_bytes(data: bytes) -> None:
    """`ooxml_guard` over the in-memory upload (same core, no temp file).

    The upload pipeline guards before the raw bytes are ever written to disk,
    which is where §2 puts the check (caps -> sniff -> OOXML guard -> dedupe).
    """
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        _guard_zip_infolist(zf.infolist())


# --- parsed representation ----------------------------------------------------

_NO_EXTRA: dict = {}  # shared read-only default; never mutated


class Block(NamedTuple):
    """One parsed unit of a document, before chunking.

    `has_table` is inherited by every chunk split out of this block — block
    level, deliberately not char-exact: marking table text inside a chunk
    would need sentinel characters, and sentinels change BM25 tokens (§2 law).
    """

    page: Optional[int]
    text: str
    has_table: bool = False
    extra: dict = _NO_EXTRA


@dataclass
class ParsedDoc:
    blocks: list[Block] = field(default_factory=list)
    tables: int = 0
    pages: Optional[int] = None


BlocksInput = Union[ParsedDoc, list]


# --- parsing ------------------------------------------------------------------

def _serialize_table(rows: list[list[Optional[str]]]) -> str:
    """Aligned `label: value` rows — financial numbers live in tables and must
    be retrievable as text."""
    lines: list[str] = []
    for row in rows or []:
        cells = [(c if c is not None else "").strip().replace("\n", " ") for c in row]
        cells = [c for c in cells if c]
        if not cells:
            continue
        if len(cells) == 1:
            lines.append(cells[0])
        else:
            lines.append(f"{cells[0]}: {' | '.join(cells[1:])}")
    return "\n".join(lines)


def _parse_pdf(path: Path, budget: Optional["_TextBudget"] = None) -> ParsedDoc:
    from pypdf import PdfReader
    import pdfplumber

    budget = budget or _TextBudget()
    reader = PdfReader(str(path))
    page_texts = []
    for page in reader.pages:  # incremental: abort on the page that crosses it
        page_texts.append(budget.keep(page.extract_text() or ""))

    table_texts: list[str] = ["" for _ in page_texts]
    table_count = 0
    try:
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages[: len(page_texts)]):
                found = [t for t in (page.extract_tables() or []) if t]
                table_count += len(found)
                serialized = [_serialize_table(t) for t in found]
                table_texts[i] = budget.keep("\n".join(s for s in serialized if s))
    except ExtractionCapExceeded:
        raise  # a cap is never "additive" — it must reach the caller
    except Exception as exc:  # noqa: BLE001 — tables are additive; page text still stands
        # Loud on purpose: this is the ONLY path where a PDF's `tables` count
        # can legitimately come back 0 for a document that has tables, and a
        # silent 0 is indistinguishable from "no tables" downstream (the TABLE
        # badge, the tables column, the upload copy).
        logger.warning(
            "pdfplumber table extraction failed (%s) — indexing page text only; "
            "this document will report tables=0 and has_table=false",
            type(exc).__name__,
        )

    blocks: list[Block] = []
    for i, (text, tables) in enumerate(zip(page_texts, table_texts), start=1):
        combined = text.strip()
        if tables:
            combined = (combined + "\n" + tables).strip()
        blocks.append(Block(i, combined, bool(tables)))
    return ParsedDoc(blocks=blocks, tables=table_count, pages=len(page_texts))


def _parse_docx(path: Path, budget: Optional["_TextBudget"] = None) -> ParsedDoc:
    import docx

    budget = budget or _TextBudget()
    document = docx.Document(str(path))
    parts = []
    for para in document.paragraphs:  # incremental: per paragraph, not per doc
        if para.text and para.text.strip():
            parts.append(budget.keep(para.text))
    has_table = False
    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            budget.add(sum(len(c or "") for c in cells))
            rows.append(cells)
        serialized = _serialize_table(rows)
        if serialized:
            parts.append(serialized)
            has_table = True
    return ParsedDoc(
        blocks=[Block(None, "\n".join(parts), has_table)],
        tables=len(document.tables),
        pages=None,
    )


def _csv_row_line(row: list, header: list[str]) -> str:
    """`col: value | col: value` — the serialization CSV and XLSX share."""
    pairs = []
    for j, value in enumerate(row):
        col = header[j] if j < len(header) and header[j] else f"col{j + 1}"
        value = "" if value is None else str(value).strip()
        if value:
            pairs.append(f"{col}: {value}")
    return " | ".join(pairs)


def _parse_csv(data: bytes, budget: Optional["_TextBudget"] = None) -> ParsedDoc:
    budget = budget or _TextBudget()
    text = _decode_text(data)
    rows = list(csv.reader(io.StringIO(text)))
    if not rows:
        return ParsedDoc(blocks=[Block(None, "")], tables=0, pages=None)
    header = [h.strip() for h in rows[0]]
    body = rows[1:] if len(rows) > 1 else []
    if not body:  # header-only (or single-row) file: keep it retrievable verbatim
        return ParsedDoc(
            blocks=[Block(None, " | ".join(h for h in header if h))], tables=0, pages=None
        )
    blocks: list[Block] = []
    for start in range(0, len(body), CSV_WINDOW_ROWS):
        lines = []
        for row in body[start : start + CSV_WINDOW_ROWS]:
            line = _csv_row_line(row, header)
            if line:
                lines.append(budget.keep(line))
        blocks.append(Block(None, "\n".join(lines), True))
    return ParsedDoc(blocks=blocks, tables=1, pages=None)


def _parse_xlsx(path: Path, budget: Optional["_TextBudget"] = None) -> ParsedDoc:
    """openpyxl, read-only + values-only: one block per ~40 rows of each sheet.

    Serialization is identical to CSV (`col: value | col: value`) so a
    spreadsheet retrieves exactly like the CSV export of the same grid. Each
    block opens with a `Sheet: {name}` line and carries the sheet name as
    chunk metadata (excluded from EMBED/LLM — the metadata-neutrality law).
    """
    import openpyxl

    budget = budget or _TextBudget()
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        sheets = wb.worksheets
        if len(sheets) > XLSX_MAX_SHEETS:
            raise ExtractionCapExceeded(
                f"spreadsheet has too many sheets (cap: {XLSX_MAX_SHEETS})"
            )

        blocks: list[Block] = []
        non_empty_sheets = 0
        cells_seen = 0
        for ws in sheets:
            title = str(ws.title)
            header: list[str] = []
            window: list[str] = []
            wrote_block = False
            has_rows = False

            def flush(force: bool = False) -> None:
                nonlocal window, wrote_block
                if not window and not force:
                    return
                text = "\n".join([f"Sheet: {title}"] + window)
                blocks.append(Block(None, text, True, {"sheet": title}))
                window = []
                wrote_block = True

            for row in ws.iter_rows(values_only=True):
                row = list(row or ())
                cells_seen += len(row)
                if cells_seen > XLSX_MAX_CELLS:
                    raise ExtractionCapExceeded(
                        f"spreadsheet too large (cap: {XLSX_MAX_CELLS} cells)"
                    )
                if not any((c is not None and str(c).strip()) for c in row):
                    continue
                has_rows = True
                if not header:  # first non-empty row is the header
                    header = [
                        (str(c).strip() if c is not None else "") for c in row
                    ]
                    continue
                line = _csv_row_line(row, header)
                if line:
                    window.append(budget.keep(line))
                if len(window) >= XLSX_WINDOW_ROWS:
                    flush()
            flush()
            if has_rows and not wrote_block:
                # header-only sheet: keep the column names retrievable
                blocks.append(
                    Block(
                        None,
                        "\n".join(
                            [f"Sheet: {title}", " | ".join(h for h in header if h)]
                        ),
                        True,
                        {"sheet": title},
                    )
                )
                wrote_block = True
            if has_rows:
                non_empty_sheets += 1
        return ParsedDoc(blocks=blocks, tables=non_empty_sheets, pages=None)
    finally:
        try:
            wb.close()
        except Exception:  # noqa: BLE001 — closing a read-only workbook is best effort
            pass


def _pptx_shapes(shapes, budget: int = 64):
    """Flatten a slide's shape tree (groups included) with a bounded walk."""
    stack = list(shapes)
    seen = 0
    while stack and seen < 4096:
        shape = stack.pop(0)
        seen += 1
        if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
            if budget > 0:
                stack.extend(list(shape.shapes))
                budget -= 1
            continue
        yield shape


def _parse_pptx(path: Path, budget: Optional["_TextBudget"] = None) -> ParsedDoc:
    """python-pptx: one block per slide (shape text, then table cell text).

    `page` is the slide number, so citations read `p.3` for slide 3 and the
    manifest's `pages` is the slide count.
    """
    from pptx import Presentation

    budget = budget or _TextBudget()
    prs = Presentation(str(path))
    slides = list(prs.slides)
    if len(slides) > PPTX_MAX_SLIDES:
        raise ExtractionCapExceeded(
            f"presentation too large (cap: {PPTX_MAX_SLIDES} slides)"
        )

    blocks: list[Block] = []
    table_count = 0
    cell_budget = PPTX_MAX_TABLE_CELLS
    truncated = False
    for index, slide in enumerate(slides, start=1):
        texts: list[str] = []
        table_texts: list[str] = []
        for shape in _pptx_shapes(slide.shapes):
            try:
                if getattr(shape, "has_text_frame", False):
                    value = (shape.text_frame.text or "").strip()
                    if value:
                        texts.append(budget.keep(value))
                if getattr(shape, "has_table", False):
                    table_count += 1
                    rows: list[list[Optional[str]]] = []
                    for row in shape.table.rows:
                        cells = list(row.cells)
                        if cell_budget <= 0:
                            truncated = True
                            break
                        cell_budget -= len(cells)
                        rows.append([c.text for c in cells])
                    serialized = _serialize_table(rows)
                    if serialized:
                        table_texts.append(budget.keep(serialized))
            except ExtractionCapExceeded:
                raise  # a cap outranks "one odd shape never fails the deck"
            except Exception:  # noqa: BLE001 — one odd shape never fails the deck
                continue
        combined = "\n".join(texts + table_texts).strip()
        blocks.append(Block(index, combined, bool(table_texts), {"slide": index}))
    if truncated:
        logger.warning(
            "pptx table extraction stopped at the %d-cell cap; remaining table "
            "cells were skipped (text still indexed)",
            PPTX_MAX_TABLE_CELLS,
        )
    return ParsedDoc(blocks=blocks, tables=table_count, pages=len(slides))


class _InertTextExtractor(HTMLParser):
    """Streaming tag stripper: script/style/comment content is dropped entirely.

    The shared extracted-text budget (MAX_EXTRACTED_TEXT_CHARS) is charged in
    `handle_data`, so a text explosion aborts mid-stream. This replaces the old
    HTML-only cap at the same threshold — one cap now covers all ten formats.

    The output is inert TEXT — it is stored, indexed and rendered as text and
    never as markup anywhere in the product (§4.1). NOTE: `convert_charrefs`
    unescapes entities, so `&lt;script&gt;` is STORED as the literal characters
    `<script>`. That is safe only because §4.1 forbids every HTML sink in the
    frontend (no `dangerouslySetInnerHTML`, ever). If anyone ever renders a
    preview, snippet or answer as rich text, this text is an XSS vector —
    re-escape at that boundary, do not "fix" it here. Raises once the collected
    text passes MAX_EXTRACTED_TEXT_CHARS so a decompression-style text
    explosion cannot run the process out of memory.
    """

    _SKIP = {"script", "style"}

    def __init__(self, budget: "_TextBudget") -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.chars = 0
        self.budget = budget
        self._skip_depth = 0

    def handle_starttag(self, tag, attrs):
        if tag in self._SKIP:
            self._skip_depth += 1

    def handle_endtag(self, tag):
        if tag in self._SKIP and self._skip_depth:
            self._skip_depth -= 1

    def handle_data(self, data):
        if self._skip_depth or not data:
            return
        self.budget.add(data)  # raises the moment the shared cap is crossed
        self.parts.append(data)
        self.chars += len(data)


def _parse_html(data: bytes, budget: Optional["_TextBudget"] = None) -> ParsedDoc:
    raw = _decode_text(data)
    parser = _InertTextExtractor(budget or _TextBudget())
    try:
        for start in range(0, len(raw), 65536):
            parser.feed(raw[start : start + 65536])
        parser.close()
    except ExtractionCapExceeded:
        raise
    except Exception:  # noqa: BLE001 — malformed markup: keep whatever was read
        logger.debug("html parse ended early; using the text collected so far")
    text = " ".join(" ".join(parser.parts).split())
    return ParsedDoc(blocks=[Block(None, text, False)], tables=0, pages=None)


def _json_scalar(value: Any) -> str:
    if value is None:
        return "null"
    if value is True:
        return "true"
    if value is False:
        return "false"
    if isinstance(value, str):
        return value
    return json.dumps(value)


def _guard_json_frontier(nodes: int, pending: int, incoming: int) -> None:
    """Refuse BEFORE pushing children, so the stack can never outgrow the cap."""
    if nodes + pending + incoming > JSON_MAX_NODES:
        raise ExtractionCapExceeded(f"json too large (cap: {JSON_MAX_NODES} nodes)")


def _parse_json(data: bytes, budget: Optional["_TextBudget"] = None) -> ParsedDoc:
    """Flatten to `a.b[0].c: value` lines, 40 lines per block.

    Traversal is ITERATIVE with an explicit stack: depth is counted, never
    recursed, so a depth bomb hits the cap instead of the interpreter's C
    stack. Any root type is allowed.

    The node cap bounds the FRONTIER, not the visit count: a flat container
    would otherwise push every sibling (each with a freshly built path string)
    before a single pop, so a cap checked on pop is structurally unable to
    bound the allocation it exists to bound (round-3 security B2).
    """
    budget = budget or _TextBudget()
    payload = json.loads(_decode_text(data))  # a failure => "failed to parse file"

    lines: list[str] = []
    nodes = 0
    stack: list[tuple[str, Any, int]] = [("", payload, 1)]
    while stack:
        path, value, depth = stack.pop()
        nodes += 1
        if nodes > JSON_MAX_NODES:
            raise ExtractionCapExceeded(f"json too large (cap: {JSON_MAX_NODES} nodes)")
        if depth > JSON_MAX_DEPTH:
            raise ExtractionCapExceeded(
                f"json too deeply nested (cap: depth {JSON_MAX_DEPTH})"
            )
        if isinstance(value, dict):
            if not value:
                lines.append(budget.keep(f"{path}: {{}}" if path else "{}"))
                continue
            _guard_json_frontier(nodes, len(stack), len(value))
            for key in reversed(list(value.keys())):
                child = f"{path}.{key}" if path else str(key)
                stack.append((child, value[key], depth + 1))
        elif isinstance(value, list):
            if not value:
                lines.append(budget.keep(f"{path}: []" if path else "[]"))
                continue
            _guard_json_frontier(nodes, len(stack), len(value))
            for i in range(len(value) - 1, -1, -1):
                child = f"{path}[{i}]" if path else f"[{i}]"
                stack.append((child, value[i], depth + 1))
        else:
            rendered = _json_scalar(value)
            lines.append(budget.keep(f"{path}: {rendered}" if path else rendered))

    blocks = [
        Block(None, "\n".join(lines[start : start + JSON_WINDOW_LINES]), False)
        for start in range(0, len(lines), JSON_WINDOW_LINES)
    ] or [Block(None, "")]
    return ParsedDoc(blocks=blocks, tables=0, pages=None)


def parse_document(path: Path, ext: str) -> ParsedDoc:
    """The real parser: blocks + the per-doc `tables` count + page/slide count.

    Counting rules (frozen, §2): PDF = pdfplumber tables; DOCX = len(tables);
    CSV = 1 when the file has data rows else 0; XLSX = one per non-empty
    worksheet; PPTX = one per table shape; TXT/MD/HTML/JSON = 0.

    ONE `_TextBudget` spans the whole document, so `MAX_EXTRACTED_TEXT_CHARS`
    bounds ACCUMULATED text across every block and every sheet/slide/page, and
    aborts mid-parse rather than after the fact.
    """
    budget = _TextBudget()
    if ext == ".pdf":
        return _parse_pdf(path, budget)
    if ext == ".docx":
        return _parse_docx(path, budget)
    if ext == ".csv":
        return _parse_csv(path.read_bytes(), budget)
    if ext in (".txt", ".md"):
        # Bounded by MAX_FILE_BYTES on the way in; charged before it is kept.
        return ParsedDoc(
            blocks=[Block(None, budget.keep(_decode_text(path.read_bytes())), False)],
            tables=0,
            pages=None,
        )
    if ext == ".xlsx":
        return _parse_xlsx(path, budget)
    if ext == ".pptx":
        return _parse_pptx(path, budget)
    if ext in (".html", ".htm"):
        return _parse_html(path.read_bytes(), budget)
    if ext == ".json":
        return _parse_json(path.read_bytes(), budget)
    raise ValueError(f"unsupported extension {ext}")


def parse_file(path: Path, ext: str) -> list[tuple[Optional[int], str]]:
    """(page, text) pairs — the v1.1 signature, kept verbatim as a test seam.

    A thin wrapper over `parse_document`; richer callers use that instead.
    """
    return [(b.page, b.text) for b in parse_document(path, ext).blocks]


# --- chunking -----------------------------------------------------------------

def provenance_prefix(doc_name: str, page: Optional[int]) -> str:
    return f"[{doc_name} — p.{page}] " if page is not None else f"[{doc_name}] "


def strip_provenance(text: str, doc_name: Optional[str], page: Optional[Any]) -> str:
    """Chunk text without its `[doc — p.N]` prefix (chunk inventory, §1.8)."""
    for prefix in (f"[{doc_name} — p.{page}] ", f"[{doc_name}] "):
        if text.startswith(prefix):
            return text[len(prefix):]
    return text


# The four v1.1 metadata keys: names, values and INSERTION ORDER are frozen.
V1_METADATA_KEYS = ("doc_id", "doc_name", "page", "chunk_ix")
# Every metadata key added in v1.2 or later. BM25Retriever tokenizes
# `node.get_content(MetadataMode.EMBED)`, which renders every non-excluded
# metadata key into the sparse text — so each of these MUST be excluded from
# both EMBED and LLM metadata modes or the sparse token stream shifts and the
# 100% eval gate moves. Dense embedding uses node.text, which is why this trap
# is invisible to dense tests. (§2 metadata-neutrality law, resolution 9.)
V2_METADATA_KEYS = ("has_table", "sheet", "slide")


def _as_blocks(blocks: BlocksInput) -> list[Block]:
    """Accept a ParsedDoc, a list[Block], or the legacy list[(page, text)]."""
    if isinstance(blocks, ParsedDoc):
        return list(blocks.blocks)
    out: list[Block] = []
    for item in blocks or []:
        if isinstance(item, Block):
            out.append(item)
        else:  # legacy 2-tuple: no table flag, no extra metadata
            page, text = item
            out.append(Block(page, text, False, _NO_EXTRA))
    return out


def chunk_pages(doc_id: str, doc_name: str, blocks: BlocksInput) -> list:
    """SentenceSplitter(512/64) nodes; metadata {doc_id, doc_name, page, chunk_ix}
    (+ v1.2 keys, all EMBED/LLM-excluded); text prefixed with provenance so the
    sparse index and the LLM both see it."""
    from llama_index.core.node_parser import SentenceSplitter
    from llama_index.core.schema import TextNode

    splitter = SentenceSplitter(chunk_size=512, chunk_overlap=64)
    nodes: list[TextNode] = []
    chunk_ix = 0
    for block in _as_blocks(blocks):
        text = block.text
        if not text or not text.strip():
            continue
        extra = block.extra or _NO_EXTRA
        for piece in splitter.split_text(text):
            if not piece.strip():
                continue
            # v1.1 keys first, in their frozen order; v1.2 keys appended after.
            metadata = {
                "doc_id": doc_id,
                "doc_name": doc_name,
                "page": block.page,
                "chunk_ix": chunk_ix,
                "has_table": bool(block.has_table),
            }
            for key in ("sheet", "slide"):
                if extra.get(key) is not None:
                    metadata[key] = extra[key]
            excluded = [k for k in V2_METADATA_KEYS if k in metadata]
            nodes.append(
                TextNode(
                    text=provenance_prefix(doc_name, block.page) + piece,
                    metadata=metadata,
                    excluded_embed_metadata_keys=list(excluded),
                    excluded_llm_metadata_keys=list(excluded),
                )
            )
            chunk_ix += 1
    return nodes


# --- ingest / delete orchestration -------------------------------------------

def _failed_entry(name: str, size: int, error: str) -> dict:
    return {
        "id": None,
        "name": name,
        "size_bytes": size,
        "pages": None,
        "chunks": 0,
        "tables": 0,
        "status": "failed",
        "error": error,
    }


def _ingest_one(name: str, data: bytes) -> dict:
    """Sync pipeline for one file (runs inside asyncio.to_thread)."""
    store = stores.get_store()
    settings = config.get_settings()
    sanitized = sanitize_filename(name)
    ext = Path(sanitized).suffix.lower()
    size = len(data)

    if ext not in ALLOWED_EXTS:
        return _failed_entry(
            sanitized, size,
            f"unsupported file type {ext or '(none)'} (allowed: {' '.join(ALLOWED_EXTS)})",
        )
    if size > MAX_FILE_BYTES:
        return _failed_entry(sanitized, size, "file exceeds the 25 MB limit")
    if not sniff_ok(ext, data):
        return _failed_entry(sanitized, size, "content does not match extension")

    # Container guard BEFORE any OOXML library (or the filesystem) sees the file.
    if ext in OOXML_EXTS:
        try:
            ooxml_guard_bytes(data)
        except ExtractionCapExceeded as exc:
            return _failed_entry(sanitized, size, exc.message)
        except Exception:  # noqa: BLE001 — not a readable container at all
            logger.warning("ooxml container unreadable for uploaded file %r", sanitized)
            return _failed_entry(sanitized, size, "failed to parse file")

    sha = hashlib.sha256(data).hexdigest()
    existing = store.find_by_sha(sha)
    if existing is not None:
        return {
            "id": existing["id"],
            "name": existing["name"],
            "size_bytes": existing["size_bytes"],
            "pages": existing["pages"],
            "chunks": existing["chunks"],
            "tables": existing.get("tables", 0) or 0,
            "status": "duplicate",
        }

    # Corpus cap — duplicates never count against it (they add nothing).
    max_documents = int(settings.max_documents)
    if len(store.get_manifest()) >= max_documents:
        return _failed_entry(
            sanitized, size,
            f"corpus is full ({max_documents} documents) — delete a document first",
        )

    doc_id = str(uuid.uuid4())
    doc_dir = config.UPLOADS_DIR / doc_id
    doc_dir.mkdir(parents=True, exist_ok=True)
    raw_path = doc_dir / sanitized
    raw_path.write_bytes(data)

    # The manifest write inside store.add_document is the ONLY commit point.
    # Every other exit — parse failure, cap exceeded, no text, provider/rate-limit
    # error, crash — must leave no trace on disk: `finally` removes the upload
    # dir whenever we did not commit (spec §1.3: failed uploads persist nothing).
    committed = False
    try:
        try:
            parsed = parse_document(raw_path, ext)
        except ExtractionCapExceeded as exc:
            # A cap is a clean per-file refusal, never a 500 and never a crash.
            logger.info("extraction cap hit for %r: %s", sanitized, exc.message)
            return _failed_entry(sanitized, size, exc.message)
        except Exception as exc:  # noqa: BLE001 — parse errors are per-file failures
            # Filename-only at WARNING (never absolute paths); full detail at DEBUG.
            logger.warning(
                "parse failed for uploaded file %r (%s)", sanitized, type(exc).__name__
            )
            logger.debug("parse failure detail for %r", sanitized, exc_info=True)
            return _failed_entry(sanitized, size, "failed to parse file")

        if not any((b.text or "").strip() for b in parsed.blocks):
            return _failed_entry(sanitized, size, "no extractable text")

        nodes = chunk_pages(doc_id, sanitized, parsed)
        if not nodes:
            return _failed_entry(sanitized, size, "no extractable text")

        vectors: Optional[list[list[float]]] = None
        bundle = providers.get_bundle()
        if bundle.provider == "gemini":
            # RateLimitedError propagates: this file + the rest of the batch fail.
            vectors = providers.embed_texts_cached(
                [n.text for n in nodes], bundle.embed_model_name
            )

        pages_count = parsed.pages
        tables_count = int(parsed.tables or 0)
        manifest_entry = {
            "id": doc_id,
            "name": sanitized,
            "ext": ext,
            "size_bytes": size,
            "sha256": sha,
            "pages": pages_count,
            "chunks": len(nodes),
            "tables": tables_count,
            "uploaded_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
            "status": "indexed",
        }
        store.add_document(manifest_entry, nodes, vectors)  # manifest write = commit
        committed = True
        return {
            "id": doc_id,
            "name": sanitized,
            "size_bytes": size,
            "pages": pages_count,
            "chunks": len(nodes),
            "tables": tables_count,
            "status": "indexed",
        }
    except providers.ProviderError:
        return _failed_entry(sanitized, size, "embedding failed (provider error)")
    finally:
        if not committed:
            shutil.rmtree(doc_dir, ignore_errors=True)


class Upload(NamedTuple):
    """A file the client sent, whose bytes have NOT been read yet.

    `read()` pulls at most `MAX_FILE_BYTES + 1` bytes when — and only when —
    the file is about to be ingested, so a request never holds more than one
    file in memory (round-3 security B3). `size` is the transport's declared
    size (an `int` or `None`); it is used only to reject early and is never
    trusted as content.
    """

    name: str
    size: Optional[int]
    read: Any  # async () -> bytes


def precheck_upload(name: str, size: Optional[int]) -> Optional[dict]:
    """Reject from NAME and DECLARED SIZE alone, before a byte is buffered.

    Returns a §1.3 `failed` entry, or None when the file must actually be read
    to decide. An upload rejected on its extension must cost the server
    nothing: previously every file was buffered first and rejected after, so
    20 x 25 MB of junk that failed the very first check still cost 494 MB.
    """
    sanitized = sanitize_filename(name)
    ext = Path(sanitized).suffix.lower()
    declared = int(size) if isinstance(size, int) and size >= 0 else 0
    if ext not in ALLOWED_EXTS:
        return _failed_entry(
            sanitized, declared,
            f"unsupported file type {ext or '(none)'} (allowed: {' '.join(ALLOWED_EXTS)})",
        )
    if declared > MAX_FILE_BYTES:
        return _failed_entry(sanitized, declared, "file exceeds the 25 MB limit")
    return None


async def _bytes_of(item) -> tuple[str, bytes]:
    """Normalize one batch item to (name, bytes), reading lazily for Uploads."""
    if isinstance(item, Upload):
        return item.name, await item.read()
    name, data = item  # legacy (name, bytes) tuple — the v1.1/v1.2 test seam
    return name, data


async def ingest_files(uploads: list) -> list[dict]:
    """Entries in upload order; rate-limit mid-batch fails the current and all
    remaining files while already-committed ones stay indexed (HTTP stays 200).

    Accepts either the frozen `list[tuple[str, bytes]]` or a list of `Upload`
    handles. With handles, a file's bytes are read immediately before its own
    ingest and released immediately after, so peak memory is ONE file — not
    the whole request.
    """
    results: list[dict] = []
    async with _ingest_lock:
        rate_limited: Optional[providers.RateLimitedError] = None
        for item in uploads:
            if rate_limited is not None:
                declared = item.size if isinstance(item, Upload) else len(item[1])
                results.append(
                    _failed_entry(
                        sanitize_filename(item.name if isinstance(item, Upload) else item[0]),
                        int(declared or 0),
                        f"rate_limited: retry in ~{rate_limited.retry_after_s}s",
                    )
                )
                continue
            # Cheap rejections happen BEFORE the bytes are pulled into RAM.
            if isinstance(item, Upload):
                early = precheck_upload(item.name, item.size)
                if early is not None:
                    results.append(early)
                    continue
            data = None
            try:
                name, data = await _bytes_of(item)
                results.append(await asyncio.to_thread(_ingest_one, name, data))
            except providers.RateLimitedError as exc:
                rate_limited = exc
                results.append(
                    _failed_entry(
                        sanitize_filename(item.name if isinstance(item, Upload) else item[0]),
                        len(data) if data is not None else 0,
                        f"rate_limited: retry in ~{exc.retry_after_s}s",
                    )
                )
            finally:
                del data  # release this file before touching the next one
    return results


async def delete_document(doc_id: str) -> bool:
    """False (=> API 404) unless doc_id is a well-formed uuid4 present in the
    manifest; uuid validation happens BEFORE any store/filesystem access."""
    if not UUID4_RE.match(doc_id or ""):
        return False
    async with _ingest_lock:
        store = stores.get_store()
        if store.find_by_id(doc_id) is None:
            return False
        await asyncio.to_thread(store.delete_document, doc_id)
        return True


# --- chunk inventory (§1.8) ---------------------------------------------------

def _preview(text: str) -> str:
    """Chunk HEAD, <=200 chars, whitespace collapsed, word-boundary cut.

    Deliberately not the citation-snippet algorithm: there is no question
    here, so there is no relevant window to centre on.
    """
    collapsed = " ".join((text or "").split())
    if len(collapsed) <= CHUNK_PREVIEW_MAX_CHARS:
        return collapsed
    piece = collapsed[: CHUNK_PREVIEW_MAX_CHARS - 1]
    cut = piece.rfind(" ")
    if cut > 0:
        piece = piece[:cut]
    return piece.rstrip() + "…"


async def chunk_inventory(doc_id: str) -> Optional[list[dict]]:
    """§1.8 rows for one document, or None when the id is malformed/unknown.

    A pure read of the docstore nodes ingest already committed: zero LLM
    calls, zero embedding calls, zero re-parsing.
    """
    if not UUID4_RE.match(doc_id or ""):
        return None
    store = stores.get_store()
    if store.find_by_id(doc_id) is None:
        return None
    rows: list[dict] = []
    for ix, node in enumerate(store.nodes_for([doc_id])):
        md = node.metadata or {}
        page = md.get("page")
        text = strip_provenance(node.get_content(), md.get("doc_name"), page)
        chunk_ix = md.get("chunk_ix")
        rows.append(
            {
                "chunk_ix": int(chunk_ix) if isinstance(chunk_ix, int) else ix,
                "page": page if isinstance(page, int) else None,
                "chars": len(text),
                # Pre-v1.2 nodes carry no has_table: absent reads as False,
                # which is tolerated state, not corruption (§3.4).
                "has_table": bool(md.get("has_table", False)),
                "preview": _preview(text),
            }
        )
    return rows


# --- startup auto-seed (§2, §5 AUTO_SEED) ------------------------------------

async def seed_sample_data() -> int:
    """Index `backend/sample_data/` when the manifest is empty. Never blocks boot.

    Goes through the normal ingest path (same lock, same commit ordering), so
    keyless mode simply indexes without vectors and the §3.5 backfill catches
    up once a key appears.
    """
    try:
        settings = config.get_settings()
        if settings.auto_seed != "on":
            return 0
        store = stores.get_store()
        if store.get_manifest():
            return 0
        source = config.SAMPLE_DATA_DIR
        if not Path(source).is_dir():
            logger.info("auto-seed skipped: no sample_data directory")
            return 0
        files = sorted(
            p
            for p in Path(source).iterdir()
            if p.is_file() and p.suffix.lower() in ALLOWED_EXTS
        )
        if not files:
            logger.info("auto-seed skipped: sample_data holds no supported files")
            return 0
        uploads = [(p.name, p.read_bytes()) for p in files]
        results = await ingest_files(uploads)
        seeded = sum(1 for r in results if r.get("status") == "indexed")
        failed = [r for r in results if r.get("status") == "failed"]
        logger.info("auto-seeded %d documents from sample_data", seeded)
        for entry in failed:
            logger.warning(
                "auto-seed could not index %r: %s", entry.get("name"), entry.get("error")
            )
        return seeded
    except Exception as exc:  # noqa: BLE001 — seeding must never block startup
        logger.warning("auto-seed failed (%s) — starting with an empty corpus", type(exc).__name__)
        return 0


# --- startup backfill (§3.5) --------------------------------------------------

def backfill_missing_embeddings() -> tuple[int, int]:
    """Gemini mode: embed docstore nodes for any manifest doc with Chroma count
    0 (indexed keyless earlier). Cache-first, batched. Rate-limit => warn and
    stop; the doc stays at count 0 (BM25 still covers it) until next restart."""
    bundle = providers.get_bundle()
    if bundle.provider != "gemini":
        return 0, 0
    store = stores.get_store()
    docs = chunks = 0
    for entry in store.get_manifest():
        if store.chroma_count_for(entry["id"]) != 0:
            continue
        nodes = store.nodes_for([entry["id"]])
        if not nodes:
            continue
        try:
            vectors = providers.embed_texts_cached(
                [n.text for n in nodes], bundle.embed_model_name
            )
        except providers.RateLimitedError as exc:
            logger.warning(
                "backfill rate-limited (retry in ~%ss) — continuing to serve; "
                "dense retrieval lacks the remaining docs until next restart",
                exc.retry_after_s,
            )
            break
        except providers.ProviderError:
            logger.warning("backfill failed (provider error) — continuing to serve")
            break
        store.add_vectors(entry["id"], nodes, vectors)
        docs += 1
        chunks += len(nodes)
    if docs:
        logger.info("backfilled %d docs / %d chunks", docs, chunks)
    return docs, chunks
